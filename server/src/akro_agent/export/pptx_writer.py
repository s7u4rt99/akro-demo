"""PPTX export: report as a presentation with clean, presentation-ready styling."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from akro_agent.models import (
    ChartSpec,
    ContentSlideSpec,
    ResearchReport,
    SlideDeckSpec,
)

# 16:9 widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color palette – professional, readable
COLOR_TITLE = RGBColor(0x1e, 0x3a, 0x5f)       # navy
COLOR_BODY = RGBColor(0x33, 0x41, 0x55)        # slate
COLOR_ACCENT = RGBColor(0x25, 0x6b, 0xeb)      # blue
COLOR_SUBTITLE = RGBColor(0x64, 0x74, 0x8b)    # gray
COLOR_WHITE = RGBColor(0xff, 0xff, 0xff)
COLOR_ACCENT_BAR = RGBColor(0x25, 0x6b, 0xeb)

# Layout indices for default blank presentation
SLIDE_LAYOUT_TITLE = 0
SLIDE_LAYOUT_TITLE_AND_BODY = 1


def _truncate_for_slide(text: str, max_chars: int = 500) -> str:
    """Truncate with ellipsis to avoid overflowing slides."""
    text = (text or "").strip()
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 3].rstrip() + "..."


def _set_title_style(shape, font_size: int = 28, color=None):
    """Apply consistent title styling to a title shape."""
    if color is None:
        color = COLOR_TITLE
    tf = shape.text_frame
    tf.word_wrap = True
    for para in tf.paragraphs:
        para.font.size = Pt(font_size)
        para.font.bold = True
        para.font.color.rgb = color
        para.space_after = Pt(4)


def _set_body_style(text_frame, font_size: int = 14, color=None):
    """Apply consistent body styling and margins."""
    if color is None:
        color = COLOR_BODY
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_right = Inches(0.3)
    text_frame.margin_top = Inches(0.15)
    text_frame.margin_bottom = Inches(0.2)
    for para in text_frame.paragraphs:
        para.font.size = Pt(font_size)
        para.font.color.rgb = color
        para.space_after = Pt(8)
        para.line_spacing = 1.15


def _add_accent_bar(slide, height: float = 0.08):
    """Add a thin accent bar under the title area."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.42),
        Inches(12.333), Inches(height),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT_BAR
    bar.line.fill.background()  # no border


def _one_point_per_slide(text: str) -> list[str]:
    """
    Split content so each point or bullet gets its own slide. No truncation.
    - Split by double newline (paragraphs); each paragraph is one slide.
    - If a paragraph has multiple lines (bullet-like), split by newline so each line = one slide.
    """
    text = (text or "").strip()
    if not text:
        return [""]
    points = []
    for block in text.split("\n\n"):
        block = block.strip()
        if not block:
            continue
        lines = [ln.strip() for ln in block.split("\n") if ln.strip()]
        if len(lines) <= 1:
            points.append(block)
        else:
            for line in lines:
                points.append(line)
    return points if points else [""]


def write_pptx(report: ResearchReport, path: str | Path) -> None:
    """Write report as a .pptx presentation with clean, presentation-ready styling."""
    path = Path(path)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ----- Title slide -----
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE])
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLOR_TITLE
    except Exception:
        pass  # some themes don't support background fill
    title_shape = slide.shapes.title
    title_shape.text = _truncate_for_slide(report.query, max_chars=80)
    for para in title_shape.text_frame.paragraphs:
        para.font.size = Pt(36)
        para.font.bold = True
        para.font.color.rgb = COLOR_WHITE
        para.alignment = 1  # center
    if len(slide.placeholders) > 1:
        sub = slide.placeholders[1].text_frame
        sub.text = "Research Report"
        for para in sub.paragraphs:
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)  # light gray
            para.alignment = 1

    # ----- Summary: one slide per point (no truncation) -----
    summary_points = _one_point_per_slide(report.summary)
    for idx, point in enumerate(summary_points):
        slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE_AND_BODY])
        slide.shapes.title.text = "Summary" if len(summary_points) == 1 else f"Summary ({idx + 1}/{len(summary_points)})"
        _set_title_style(slide.shapes.title, font_size=24)
        _add_accent_bar(slide)
        tf = slide.placeholders[1].text_frame
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_BODY
        p.space_after = Pt(10)
        p.line_spacing = 1.2
        _set_body_style(tf, font_size=15)

    # ----- Section slides: one slide per point/bullet (no truncation) -----
    for sec in report.sections:
        title = sec.get("title") or "Section"
        content = sec.get("content") or ""
        points = _one_point_per_slide(content)
        for idx, point in enumerate(points):
            slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE_AND_BODY])
            slide_title = title if len(points) == 1 else f"{title} ({idx + 1}/{len(points)})"
            slide.shapes.title.text = slide_title[:75] if len(slide_title) > 75 else slide_title
            _set_title_style(slide.shapes.title, font_size=22)
            _add_accent_bar(slide)
            tf = slide.placeholders[1].text_frame
            p = tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_BODY
            p.space_after = Pt(8)
            p.line_spacing = 1.15
            _set_body_style(tf, font_size=13)

    # ----- Confidence notes: one slide per point (no truncation) -----
    if report.confidence_notes:
        conf_points = _one_point_per_slide(report.confidence_notes)
        for idx, point in enumerate(conf_points):
            slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE_AND_BODY])
            slide.shapes.title.text = "Confidence & limitations" if len(conf_points) == 1 else f"Confidence & limitations ({idx + 1}/{len(conf_points)})"
            _set_title_style(slide.shapes.title, font_size=22)
            _add_accent_bar(slide)
            tf = slide.placeholders[1].text_frame
            p = tf.paragraphs[0]
            p.text = point
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_BODY
            p.space_after = Pt(8)
            _set_body_style(tf, font_size=13)

    # ----- Sources: one slide per source (no truncation) -----
    if report.sources:
        for idx, src in enumerate(report.sources):
            slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE_AND_BODY])
            slide.shapes.title.text = "Sources" if len(report.sources) == 1 else f"Sources ({idx + 1}/{len(report.sources)})"
            _set_title_style(slide.shapes.title, font_size=22)
            _add_accent_bar(slide)
            tf = slide.placeholders[1].text_frame
            p = tf.paragraphs[0]
            p.text = src
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_SUBTITLE
            p.word_wrap = True
            _set_body_style(tf, font_size=10)

    prs.save(str(path))


def _add_chart_to_slide(slide, chart_spec: ChartSpec, left: float, top: float, width: float, height: float) -> None:
    """Add a chart to the slide from ChartSpec."""
    chart_type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
    }
    xl_type = chart_type_map.get(chart_spec.chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    categories = chart_spec.categories or ["A", "B", "C"][: len(chart_spec.values)]
    values = chart_spec.values or [1, 2, 3]
    if len(categories) != len(values):
        values = values[: len(categories)] or [1] * len(categories)
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(chart_spec.series_name or "Value", tuple(values))
    slide.shapes.add_chart(
        xl_type,
        Inches(left), Inches(top), Inches(width), Inches(height),
        chart_data,
    )


def write_pptx_from_spec(spec: SlideDeckSpec, path: str | Path) -> None:
    """Build a .pptx from an AI-generated slide deck spec (icons, bullets, optional charts)."""
    path = Path(path)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ----- Title slide -----
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE])
    try:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COLOR_TITLE
    except Exception:
        pass
    title_shape = slide.shapes.title
    title_shape.text = _truncate_for_slide(spec.title_slide_title, max_chars=80)
    for para in title_shape.text_frame.paragraphs:
        para.font.size = Pt(36)
        para.font.bold = True
        para.font.color.rgb = COLOR_WHITE
        para.alignment = 1
    if len(slide.placeholders) > 1:
        sub = slide.placeholders[1].text_frame
        sub.text = spec.title_slide_subtitle or "Research Report"
        for para in sub.paragraphs:
            para.font.size = Pt(18)
            para.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
            para.alignment = 1

    # ----- Content slides (from spec) -----
    for content in spec.content_slides:
        slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_TITLE_AND_BODY])
        # Title with icon (emoji) for visual interest
        title_with_icon = f"{content.icon}  {_truncate_for_slide(content.title, max_chars=70)}"
        slide.shapes.title.text = title_with_icon
        _set_title_style(slide.shapes.title, font_size=22)
        _add_accent_bar(slide)

        bullets = content.bullets or []
        has_chart = content.chart and content.chart.categories and content.chart.values
        tf = slide.placeholders[1].text_frame
        for i, bullet in enumerate(bullets[:5] if has_chart else bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(12) if has_chart else Pt(14)
            p.font.color.rgb = COLOR_BODY
            p.space_after = Pt(4) if has_chart else Pt(6)
        _set_body_style(tf, font_size=12 if has_chart else 14)
        if has_chart:
            _add_chart_to_slide(
                slide, content.chart,
                left=1.0, top=3.8, width=11.0, height=3.2,
            )

    prs.save(str(path))
